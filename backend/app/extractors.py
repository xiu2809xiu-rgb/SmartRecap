from io import BytesIO
from pathlib import Path
import re
from typing import List, Tuple
from zipfile import BadZipFile, ZipFile

import fitz
from docx import Document
from pptx import Presentation


class ExtractionError(ValueError):
    pass


ALLOWED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".txt", ".md", ".png", ".jpg", ".jpeg"}


def validate_file(content: bytes, filename: str, max_bytes: int) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix not in ALLOWED_EXTENSIONS:
        raise ExtractionError("Unsupported format. Use PDF, PPTX, DOCX, TXT, MD, PNG, or JPG.")
    if not content:
        raise ExtractionError("The uploaded file is empty.")
    if len(content) > max_bytes:
        raise ExtractionError("The uploaded file exceeds the configured size limit.")
    if suffix == ".pdf" and not content.startswith(b"%PDF"):
        raise ExtractionError("The file extension says PDF, but the content is not a valid PDF.")
    if suffix in {".pptx", ".docx"}:
        _validate_office_archive(content, suffix)
    if suffix == ".png" and not content.startswith(b"\x89PNG\r\n\x1a\n"):
        raise ExtractionError("Invalid PNG signature.")
    if suffix in {".jpg", ".jpeg"} and not content.startswith(b"\xff\xd8\xff"):
        raise ExtractionError("Invalid JPEG signature.")
    return suffix


def _validate_office_archive(content: bytes, suffix: str) -> None:
    try:
        with ZipFile(BytesIO(content)) as archive:
            names = set(archive.namelist())
    except BadZipFile as exc:
        raise ExtractionError("The Office document is not a valid archive.") from exc
    marker = "word/document.xml" if suffix == ".docx" else "ppt/presentation.xml"
    if marker not in names:
        raise ExtractionError("The file content does not match its Office document extension.")


def extract_locally(content: bytes, suffix: str) -> Tuple[str, List[str]]:
    if suffix == ".pdf":
        return _extract_pdf(content)
    if suffix == ".pptx":
        return _extract_pptx(content)
    if suffix == ".docx":
        return _extract_docx(content)
    if suffix in {".txt", ".md"}:
        text = content.decode("utf-8", errors="replace").strip()
        return "[Section 1]\n" + text, ["Section 1"]
    raise ExtractionError("This image requires Azure Content Understanding OCR.")


def _extract_pdf(content: bytes) -> Tuple[str, List[str]]:
    pages, labels = [], []
    try:
        document = fitz.open(stream=content, filetype="pdf")
        if document.page_count > 250:
            raise ExtractionError("PDFs are limited to 250 pages.")
        for index, page in enumerate(document):
            label = "Page {}".format(index + 1)
            text = page.get_text("text", sort=True).strip()
            if text:
                pages.append("[{}]\n{}".format(label, text))
                labels.append(label)
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError("The PDF could not be read.") from exc
    return "\n\n".join(pages), labels

def _extract_pptx(content: bytes) -> Tuple[str, List[str]]:
    slides, labels = [], []
    try:
        presentation = Presentation(BytesIO(content))
        if len(presentation.slides) > 300:
            raise ExtractionError("Presentations are limited to 300 slides.")
        for index, slide in enumerate(presentation.slides):
            fragments = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    fragments.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        fragments.append(" | ".join(cell.text.strip() for cell in row.cells))
            if fragments:
                label = "Slide {}".format(index + 1)
                slides.append("[{}]\n{}".format(label, "\n".join(fragments)))
                labels.append(label)
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError("The PowerPoint file could not be read.") from exc
    return "\n\n".join(slides), labels


def _extract_docx(content: bytes) -> Tuple[str, List[str]]:
    sections, labels = [], []
    try:
        document = Document(BytesIO(content))
        current, number = [], 1
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            if paragraph.style and paragraph.style.name.startswith("Heading") and current:
                label = "Section {}".format(number)
                sections.append("[{}]\n{}".format(label, "\n".join(current)))
                labels.append(label)
                number += 1
                current = [text]
            else:
                current.append(text)
        for table in document.tables:
            for row in table.rows:
                current.append(" | ".join(cell.text.strip() for cell in row.cells))
        if current:
            label = "Section {}".format(number)
            sections.append("[{}]\n{}".format(label, "\n".join(current)))
            labels.append(label)
    except Exception as exc:
        raise ExtractionError("The Word document could not be read.") from exc
    return "\n\n".join(sections), labels


def extract_with_azure(content: bytes, content_type: str, endpoint: str, key: str, analyzer_id: str) -> Tuple[str, List[str]]:
    """Use Content Understanding for scanned documents and images."""
    try:
        from azure.ai.contentunderstanding import ContentUnderstandingClient
        from azure.core.credentials import AzureKeyCredential

        client = ContentUnderstandingClient(endpoint=endpoint.rstrip("/"), credential=AzureKeyCredential(key))
        poller = client.begin_analyze(analyzer_id=analyzer_id, body=BytesIO(content), content_type=content_type)
        result = poller.result()
        payload = result.as_dict() if hasattr(result, "as_dict") else result
        text = _find_longest_text(payload)
        if not text.strip():
            raise ExtractionError("Azure OCR returned no readable text.")
        return "[Page 1]\n" + text.strip(), ["Page 1"]
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError("Azure Content Understanding could not analyze this file.") from exc


def _find_longest_text(value) -> str:
    candidates = []
    if isinstance(value, str):
        return value
    if isinstance(value, dict):
        for key, child in value.items():
            found = _find_longest_text(child)
            if found and key.lower() in {"markdown", "content", "text"}:
                candidates.append(found)
            elif found:
                candidates.append(found)
    elif isinstance(value, list):
        candidates.extend(_find_longest_text(child) for child in value)
    return max((item for item in candidates if item), key=len, default="")


from threading import Lock

_OCR_ENGINE = None
_PADDLE_ENGINE = None
_MATH_ENGINE = None
_RAPID_LOCK = Lock()
_PADDLE_LOCK = Lock()
_MATH_LOCK = Lock()


def paddle_ocr_available() -> bool:
    import sys

    # PaddleOCR 3.x imports PaddleX/ModelScope, which loads PyTorch after Paddle.
    # Their Windows DLL runtimes conflict in the combined OCR environment; EC2
    # runs Linux, where the pinned CPU stack remains available.
    if sys.platform == "win32":
        return False
    try:
        from importlib.util import find_spec
        return find_spec("paddle") is not None and find_spec("paddleocr") is not None
    except (ImportError, ValueError):
        return False


def math_ocr_available() -> bool:
    try:
        from importlib.util import find_spec
        return find_spec("pix2text") is not None
    except (ImportError, ValueError):
        return False


def _corruption_score(text: str) -> int:
    return len(re.findall(r"(?:�|\(cid:\d+\)|[□▯◻�]|[\ue000-\uf8ff])", text))


def _math_page_score(text: str) -> int:
    operators = "∫∑√∞≈≠≤≥→∂∇×·±⊥∥∈∉∪∩"
    score = sum(text.count(symbol) * 2 for symbol in operators)
    score += len(re.findall(r"(?:\^|_)[A-Za-z0-9{]", text))
    score += len(re.findall(r"\b(?:sin|cos|tan|lim|log|ln|det|grad|div)\b", text, flags=re.IGNORECASE))
    score += len(re.findall(r"[A-Za-z0-9)]\s*[=+*/]\s*[A-Za-z0-9(]", text))
    score += 5 * _corruption_score(text)
    if text and len(re.findall(r"\b[A-Za-z]\b", text)) >= 8:
        score += 5
    return score


def _layout_math_score(page) -> int:
    """Detect raised/lowered small PDF spans that plain-text extraction flattens."""
    try:
        spans = [
            span
            for block in page.get_text("dict", sort=True).get("blocks", [])
            for line in block.get("lines", [])
            for span in line.get("spans", [])
            if str(span.get("text") or "").strip() and float(span.get("size") or 0) > 0
        ]
        sizes = sorted(float(span["size"]) for span in spans)
        if not sizes:
            return 0
        body_size = sizes[len(sizes) // 2]
        return sum(
            2
            for span in spans
            if float(span["size"]) <= body_size * 0.82
            and len(str(span.get("text") or "").strip()) <= 12
            and re.search(r"[A-Za-z0-9+−=]", str(span.get("text") or ""))
        )
    except Exception:
        return 0


def _labeled_sections(text: str) -> List[Tuple[str, str]]:
    matches = list(re.finditer(r"(?m)^\[([^\]\n]+)\]\s*$", text or ""))
    if not matches:
        return [("Section 1", (text or "").strip())] if (text or "").strip() else []
    return [
        (
            match.group(1).strip(),
            (text[match.end():matches[index + 1].start()] if index + 1 < len(matches) else text[match.end():]).strip(),
        )
        for index, match in enumerate(matches)
    ]


def _math_structure_score(text: str) -> int:
    return len(re.findall(r"(?:\$|\\\[|\\\(|\\frac|\\sqrt|\^\{|_\{|[∫∑√])", text or ""))


def merge_extracted_text(base: str, recovered: str, math: bool = False) -> str:
    """Merge locator sections without allowing recovered pages to be shadowed."""
    sections = {label: value for label, value in _labeled_sections(base)}
    order = [label for label, _ in _labeled_sections(base)]
    for label, value in _labeled_sections(recovered):
        if not value:
            continue
        existing = sections.get(label, "")
        if label not in sections:
            order.append(label)
            sections[label] = value
            continue
        replace = (
            _corruption_score(existing) > _corruption_score(value)
            or (len(existing.strip()) < 80 and len(value.strip()) >= len(existing.strip()))
            or (
                math
                and _math_structure_score(value) >= max(3, _math_structure_score(existing) + 2)
            )
        )
        if replace:
            sections[label] = value
    return "\n\n".join("[{}]\n{}".format(label, sections[label]) for label in order if sections.get(label))


def extract_with_math_ocr(
    content: bytes,
    suffix: str,
    max_pages: int = 8,
    time_budget_seconds: int = 120,
) -> Tuple[str, List[str]]:
    """Recognize selected formula-heavy PDF pages as Markdown/LaTeX on CPU."""
    if suffix != ".pdf" or not math_ocr_available():
        return "", []
    import time

    deadline = time.monotonic() + max(10, time_budget_seconds)
    sections, labels = [], []
    try:
        from PIL import Image
        from pix2text import Pix2Text
        from tempfile import TemporaryDirectory

        with fitz.open(stream=content, filetype="pdf") as document:
            ranked = []
            for index, page in enumerate(document):
                native_text = page.get_text("text", sort=True).strip()
                score = _math_page_score(native_text) + _layout_math_score(page)
                if score >= 8 or len(native_text) < 80:
                    ranked.append((score, index))
            limit = max(1, min(max_pages, 12))
            candidates = sorted(index for _, index in sorted(ranked, key=lambda item: (-item[0], item[1]))[:limit])
            if not candidates:
                return "", []
            global _MATH_ENGINE
            with _MATH_LOCK:
                if _MATH_ENGINE is None:
                    _MATH_ENGINE = Pix2Text.from_config(
                        enable_formula=True,
                        enable_table=False,
                        device="cpu",
                    )
            for index in candidates:
                if time.monotonic() >= deadline:
                    break
                page = document[index]
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
                image = Image.open(BytesIO(pixmap.tobytes("png"))).convert("RGB")
                result = _MATH_ENGINE.recognize_page(
                    image,
                    page_number=index,
                    resized_shape=1280,
                    mfr_batch_size=1,
                )
                with TemporaryDirectory(prefix="smartrecap-math-") as output_dir:
                    markdown = result.to_markdown(output_dir, markdown_fn=None) if hasattr(result, "to_markdown") else str(result)
                markdown = str(markdown or "").strip()
                if markdown:
                    label = "Page {}".format(index + 1)
                    sections.append("[{}]\n{}".format(label, markdown))
                    labels.append(label)
    except Exception as exc:
        raise ExtractionError("Pix2Text math OCR could not process the selected formula pages.") from exc
    return "\n\n".join(sections), labels


def local_ocr_available() -> bool:
    try:
        from importlib.util import find_spec
        return find_spec("rapidocr_onnxruntime") is not None or paddle_ocr_available()
    except (ImportError, ValueError):
        return False


def extract_with_local_ocr(content: bytes, suffix: str, native_text: str = "", deep_scan: bool = False, max_images: int = 40, time_budget_seconds: int = 60, use_paddle: bool = True) -> Tuple[str, List[str]]:
    """Scan only image-dependent content, with strict image and time budgets."""
    import time

    images = []
    limit = max(1, min(max_images, 24 if deep_scan else 40))
    deadline = time.monotonic() + max(5, time_budget_seconds)
    try:
        if suffix in {".png", ".jpg", ".jpeg"}:
            images.append(("Image 1", content))
        elif suffix == ".pdf":
            with fitz.open(stream=content, filetype="pdf") as document:
                candidates = [
                    index
                    for index, page in enumerate(document)
                    if len(page.get_text("text").strip()) < 80
                ]
                if len(candidates) > limit:
                    step = (len(candidates) - 1) / max(1, limit - 1)
                    candidates = list(dict.fromkeys(candidates[round(position * step)] for position in range(limit)))
                for index in candidates:
                    if time.monotonic() >= deadline:
                        break
                    page = document[index]
                    pixmap = page.get_pixmap(matrix=fitz.Matrix(1.35, 1.35), alpha=False)
                    images.append(("Page {}".format(index + 1), pixmap.tobytes("png")))
        elif suffix == ".pptx":
            presentation = Presentation(BytesIO(content))
            for slide_index, slide in enumerate(presentation.slides):
                if len(images) >= limit or time.monotonic() >= deadline:
                    break
                slide_text = " ".join(shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
                if deep_scan or len(slide_text) < 100:
                    for shape in slide.shapes:
                        if len(images) >= limit:
                            break
                        if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                            images.append(("Slide {} image".format(slide_index + 1), shape.image.blob))
        elif suffix == ".docx" and (deep_scan or len(native_text.strip()) < 800):
            with ZipFile(BytesIO(content)) as archive:
                media = [name for name in archive.namelist() if name.startswith("word/media/")][:limit]
                images.extend(("Embedded image {}".format(index + 1), archive.read(name)) for index, name in enumerate(media))
    except Exception as exc:
        raise ExtractionError("Local OCR could not prepare this document for scanning.") from exc
    sections, labels = [], []
    for label, image in images[:limit]:
        if time.monotonic() >= deadline:
            break
        text = _ocr_image(image, use_paddle=use_paddle, prefer_paddle=deep_scan)
        if text:
            sections.append("[{}]\n{}".format(label, text))
            labels.append(label)
    return "\n\n".join(sections), labels


def _prepare_image(content: bytes):
    try:
        import numpy as np
        from PIL import Image, UnidentifiedImageError
    except ImportError as exc:
        raise ExtractionError("Local OCR is unavailable. Install the backend OCR dependencies.") from exc
    try:
        image = Image.open(BytesIO(content)).convert("RGB")
    except (UnidentifiedImageError, OSError, ValueError):
        return None
    if image.width * image.height > 28_000_000:
        image.thumbnail((4200, 4200))
    return np.asarray(image)


def _rapid_ocr(image) -> Tuple[str, float]:
    global _OCR_ENGINE
    try:
        from rapidocr_onnxruntime import RapidOCR
    except ImportError:
        return "", 0.0
    try:
        with _RAPID_LOCK:
            if _OCR_ENGINE is None:
                _OCR_ENGINE = RapidOCR()
            result, _ = _OCR_ENGINE(image)
    except Exception as exc:
        raise ExtractionError("RapidOCR could not read an image in this source.") from exc
    if not result:
        return "", 0.0
    rows = [(str(item[1]).strip(), float(item[2]) if len(item) > 2 else 0.0) for item in result if len(item) > 1 and str(item[1]).strip()]
    if not rows:
        return "", 0.0
    return "\n".join(row[0] for row in rows), sum(row[1] for row in rows) / len(rows)


def _paddle_ocr(image) -> Tuple[str, float]:
    global _PADDLE_ENGINE
    if not paddle_ocr_available():
        return "", 0.0
    try:
        import json
        from paddleocr import PaddleOCR

        with _PADDLE_LOCK:
            if _PADDLE_ENGINE is None:
                _PADDLE_ENGINE = PaddleOCR(
                    device="cpu",
                    text_detection_model_name="PP-OCRv5_mobile_det",
                    text_recognition_model_name="en_PP-OCRv5_mobile_rec",
                    text_det_limit_side_len=1280,
                    text_det_limit_type="max",
                    enable_mkldnn=False,
                    use_doc_orientation_classify=False,
                    use_doc_unwarping=False,
                    use_textline_orientation=False,
                )
            outputs = list(_PADDLE_ENGINE.predict(input=image))
        texts, scores = [], []
        for output in outputs:
            payload = getattr(output, "json", output)
            if callable(payload):
                payload = payload()
            if isinstance(payload, str):
                payload = json.loads(payload)
            if not isinstance(payload, dict):
                continue
            result = payload.get("res", payload)
            found_texts = result.get("rec_texts", []) or []
            found_scores = result.get("rec_scores", []) or []
            texts.extend(str(text).strip() for text in found_texts if str(text).strip())
            scores.extend(float(score) for score in found_scores)
        confidence = sum(scores) / len(scores) if scores else 0.0
        return "\n".join(texts), confidence
    except Exception as exc:
        import logging
        logging.getLogger("smartrecap.ocr").warning("PaddleOCR fallback failed: %s", exc)
        return "", 0.0


def _ocr_image(content: bytes, use_paddle: bool = True, prefer_paddle: bool = False) -> str:
    image = _prepare_image(content)
    if image is None:
        return ""
    if use_paddle and prefer_paddle:
        paddle_text, paddle_confidence = _paddle_ocr(image)
        if paddle_text and paddle_confidence >= 0.55:
            return paddle_text
        rapid_text, rapid_confidence = _rapid_ocr(image)
        return paddle_text if paddle_text and paddle_confidence >= rapid_confidence else rapid_text
    rapid_text, rapid_confidence = _rapid_ocr(image)
    if use_paddle and (len(rapid_text.strip()) < 24 or rapid_confidence < 0.72):
        paddle_text, paddle_confidence = _paddle_ocr(image)
        if paddle_text and (not rapid_text or paddle_confidence >= rapid_confidence or len(paddle_text) > len(rapid_text) * 1.15):
            return paddle_text
    return rapid_text